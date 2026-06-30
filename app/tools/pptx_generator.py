import os
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation(slide_data: List[Dict[str, Any]], output_filepath: str):
    """Generates a professional PowerPoint presentation based on structured slide data.
    
    Args:
        slide_data: A list of slide dictionaries:
            [
                {"title": "Title Slide", "content": ["Subtitle Line 1", "Subtitle Line 2"], "is_title": True},
                {"title": "Slide Title", "content": ["Bullet Point 1", "Bullet Point 2"]}
            ]
        output_filepath: Target path to save the .pptx file.
    """
    prs = Presentation()
    
    # Standard widescreen (16:9) dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Professional, rich color scheme (Deep Navy, Slate Blue, Charcoal Text)
    PRIMARY_COLOR = RGBColor(15, 32, 67)      # Deep Navy
    SECONDARY_COLOR = RGBColor(79, 110, 138)  # Slate Blue
    TEXT_COLOR = RGBColor(40, 40, 40)         # Slate Charcoal
    
    for slide_info in slide_data:
        is_title = slide_info.get("is_title", False)
        title_text = slide_info.get("title", "")
        content_items = slide_info.get("content", [])
        
        if is_title:
            # Layout 0 is typically the Title layout
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            
            title_shape.text = title_text
            # Customize title font
            title_para = title_shape.text_frame.paragraphs[0]
            title_para.font.name = "Arial" # Default safe font
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.font.color.rgb = PRIMARY_COLOR
            
            # Format subtitle text
            if isinstance(content_items, list):
                subtitle_shape.text = "\n".join(content_items)
            else:
                subtitle_shape.text = str(content_items)
                
            sub_para = subtitle_shape.text_frame.paragraphs[0]
            sub_para.font.name = "Arial"
            sub_para.font.size = Pt(20)
            sub_para.font.color.rgb = SECONDARY_COLOR
        else:
            # Layout 1 is Title and Content
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            title_shape.text = title_text
            
            title_para = title_shape.text_frame.paragraphs[0]
            title_para.font.name = "Arial"
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = PRIMARY_COLOR
            
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear() # Clear default placeholder bullets
            
            for i, item in enumerate(content_items):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = item
                p.level = 0
                p.font.name = "Arial"
                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_COLOR
                p.space_after = Pt(10)
                
    # Ensure directory exists and save presentation
    os.makedirs(os.path.dirname(output_filepath), exist_ok=True)
    prs.save(output_filepath)
    print(f"PowerPoint presentation created at: {output_filepath}")
